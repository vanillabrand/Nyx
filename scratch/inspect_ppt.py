from pptx import Presentation
import os

path = r"C:\Users\bruce\Downloads\AI VENTRA PPT TEMPLATE.pptx"
if not os.path.exists(path):
    print(f"File not found: {path}")
    exit(1)

prs = Presentation(path)

print(f"Total slides: {len(prs.slides)}")
for i, slide in enumerate(prs.slides):
    print(f"\nSlide {i+1}:")
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            print(f" - [Shape]: {shape.text[:50]}...")
