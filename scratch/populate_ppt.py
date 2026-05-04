from pptx import Presentation
import os

def replace_text(slide, search_text, replacement_text):
    for shape in slide.shapes:
        if hasattr(shape, "text") and search_text in shape.text:
            shape.text = shape.text.replace(search_text, replacement_text)

path = r"C:\Users\bruce\Downloads\AI VENTRA PPT TEMPLATE.pptx"
output_path = r"C:\Users\bruce\Downloads\Aviation_Incident_Command_Presentation.pptx"

if not os.path.exists(path):
    print(f"File not found: {path}")
    exit(1)

prs = Presentation(path)

# Slide 1: Title
# We'll try to find a title shape
if len(prs.slides) > 0:
    for shape in prs.slides[0].shapes:
        if hasattr(shape, "text"):
            if "TITLE" in shape.text.upper() or "NAME" in shape.text.upper() or "TEAM" in shape.text.upper():
                shape.text = "Aviation Incident Command: Tactical HUD"

# Slide 2: Problem Understanding
if len(prs.slides) > 1:
    content = (
        "Aviation data is siloed between flight telemetry and incident reports. "
        "Analysts lack a unified, high-fidelity dashboard to merge real-time positions with narrative safety intelligence. "
        "Key Challenges: Site blocks, data fragmentation, and mapping messy text to official airport codes."
    )
    for shape in prs.slides[1].shapes:
        if hasattr(shape, "text") and "interpretation" in shape.text.lower():
            shape.text = content

# Slide 3: Solution Overview
if len(prs.slides) > 2:
    content = (
        "International Flight Incident & Path Tracking Utility. "
        "A Parallel Data Stream Engine with Neo4j and AI models trained on archived flight data to mine potential safety issues in real-time. "
        "Features: Live 3D globe, narrative hydration, and cockpit-grade UI."
    )
    for shape in prs.slides[2].shapes:
        if hasattr(shape, "text") and "solution proposal" in shape.text.lower():
            shape.text = content

# Slide 4: System Architecture
if len(prs.slides) > 3:
    content = (
        "Frontend: React/Vite/Three.js with a custom Cockpit UI. "
        "Backend: Parallel TypeScript Scraper Suite with ProxyJet (SOCKS5) and Direct-Link fallbacks. "
        "Intelligence Layer: Neo4j Graph DB and Custom Pattern Spotting AI trained on history."
    )
    for shape in prs.slides[3].shapes:
        if hasattr(shape, "text") and "Block Diagram" in shape.text:
            # We can't easily draw a diagram, so we replace text if there's a text box
            pass
        elif hasattr(shape, "text"):
            # Try to find a content area
            if len(shape.text) > 10 and "Slide" not in shape.text:
                 # shape.text = content
                 pass

# Slide 5: Challenges & Accomplishments
if len(prs.slides) > 4:
    content = (
        "Challenges: Handling aggressive scraper blocks and synchronizing data from multiple fragmented sources into a unified schema. "
        "Accomplishments: A smooth, lag-free web app design that handles live 3D globes and constant data streams on any modern device."
    )
    for shape in prs.slides[4].shapes:
        if hasattr(shape, "text"):
            if "Key findings" in shape.text or "Learning" in shape.text:
                 shape.text = content

# Slide 6: What's Next
if len(prs.slides) > 5:
    content = (
        "Predictive Mining: Using Neo4j and AI to flag recurring patterns in safety data before they become accidents. "
        "Conversational AI: Adding a natural language query engine to allow users to talk to the globe."
    )
    for shape in prs.slides[5].shapes:
        if hasattr(shape, "text"):
            if "Next steps" in shape.text or "Future" in shape.text:
                 shape.text = content

prs.save(output_path)
print(f"Presentation saved to: {output_path}")
